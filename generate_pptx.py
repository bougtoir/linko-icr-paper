#!/usr/bin/env python3
"""Build the editable English figure deck for the LINKO manuscript.

One slide per manuscript figure: editable title, the generated image, and the
editable figure legend. Figure list and legends come from
:mod:`icr_paper.src.manuscript_content`, so the deck cannot drift from the
manuscript. Run ``python run_analysis.py`` first.
"""

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
sys.path.insert(0, str(BASE.parent))

from icr_paper.src.manuscript_content import (
    build_english,
    figure_blocks,
    prepare_manuscript,
)
from icr_paper.src.results_loader import load_results

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Calibri"
TITLE = (
    "LINKO (Latent Information Normalization for Key Outcomes): A Framework for "
    "Evaluating the Validity of Meta-Analytic Pooling Across Heterogeneous RCT Data Structures"
)


def _textbox(slide, left, top, width, height, text, size, *, bold=False,
             italic=False, colour=RGBColor(0x22, 0x22, 0x22),
             align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    run.font.color.rgb = colour
    paragraph.alignment = align
    return box


def add_title_slide(prs: Presentation, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _textbox(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.2), TITLE, 30,
             bold=True)
    _textbox(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.5), subtitle, 16,
             colour=RGBColor(0x55, 0x55, 0x55))


def add_figure_slide(prs: Presentation, spec: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _textbox(slide, Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.7),
             spec["label"], 22, bold=True)

    path = Path(spec["path"])
    if not path.exists():
        raise SystemExit(f"Missing figure {path}; run run_analysis.py first.")
    with Image.open(path) as image:
        aspect = image.size[0] / image.size[1]
    max_w, max_h = Inches(12.0), Inches(5.4)
    if aspect > max_w / max_h:
        width, height = max_w, int(max_w / aspect)
    else:
        height, width = max_h, int(max_h * aspect)
    slide.shapes.add_picture(str(path), int((SLIDE_W - width) / 2), Inches(0.95),
                             width, height)

    _textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.9),
             f"{spec['label']}. {spec['caption']}", 12, italic=True,
             colour=RGBColor(0x44, 0x44, 0x44))


def main() -> None:
    results = load_results()
    blocks = prepare_manuscript(build_english(results))
    specs = figure_blocks(blocks)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_title_slide(
        prs,
        f"Editable figures ({len(specs)}), generated from "
        f"results/results.json on {results['metadata']['generated_utc']}",
    )
    for spec in specs:
        add_figure_slide(prs, spec)

    output = BASE / "LINKO_figures_english.pptx"
    prs.save(str(output))
    print(f"Wrote {output} ({len(specs)} figure slides)")


if __name__ == "__main__":
    main()
