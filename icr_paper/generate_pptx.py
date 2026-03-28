#!/usr/bin/env python3
"""Generate PPTX figure files (English + Japanese) for LINKO paper.

Each slide has one figure or table:
- Code-generated plots: embedded as images
- Flow/concept diagrams: created as editable PowerPoint shapes
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(BASE, "figures")

# Slide dimensions: widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_pres():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title box
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.0))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    tx2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.3), Inches(1.5))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.name = "Calibri"
    p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p2.alignment = PP_ALIGN.CENTER
    return slide


def add_image_slide(prs, fig_filename, title_text, caption_text=""):
    """Add a slide with a code-generated figure embedded as an image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title
    tx = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Image
    fig_path = os.path.join(FIG_DIR, fig_filename)
    if os.path.exists(fig_path):
        # Center image on slide
        from PIL import Image
        img = Image.open(fig_path)
        img_w, img_h = img.size
        aspect = img_w / img_h

        max_w = Inches(12.0)
        max_h = Inches(5.8)

        if aspect > (max_w / max_h):
            w = max_w
            h = int(w / aspect)
        else:
            h = max_h
            w = int(h * aspect)

        left = int((SLIDE_W - w) / 2)
        top = Inches(0.95)
        slide.shapes.add_picture(fig_path, left, top, w, h)
    else:
        tx2 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
        tf2 = tx2.text_frame
        tf2.paragraphs[0].text = f"[Figure not found: {fig_filename}]"
        tf2.paragraphs[0].font.size = Pt(16)
        tf2.paragraphs[0].font.italic = True

    # Caption
    if caption_text:
        tx3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.5))
        tf3 = tx3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = caption_text
        p3.font.size = Pt(11)
        p3.font.italic = True
        p3.font.name = "Calibri"
        p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p3.alignment = PP_ALIGN.CENTER

    return slide


def _add_shape_box(slide, left, top, width, height, text, fill_rgb, font_size=12, bold=False, align=PP_ALIGN.CENTER):
    """Helper to add a rounded rectangle with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].font.name = "Calibri"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def _add_arrow(slide, start_x, start_y, end_x, end_y):
    """Add a connector arrow between two points."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_x, start_y, end_x, end_y
    )
    connector.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    connector.line.width = Pt(2)
    # Add arrowhead
    from pptx.oxml.ns import qn
    ln = connector.line._ln
    tail_end = ln.makeelement(qn('a:tailEnd'), {})
    tail_end.set('type', 'triangle')
    tail_end.set('w', 'med')
    tail_end.set('len', 'med')
    ln.append(tail_end)
    return connector


def add_linko_framework_slide(prs, lang="en"):
    """Create editable LINKO Framework Overview diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_text = "LINKO Framework Overview" if lang == "en" else "LINKO\u30d5\u30ec\u30fc\u30e0\u30ef\u30fc\u30af\u6982\u8981"
    tx = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Top: LINKO title box
    linko_title = "LINKO\n(Latent Information Normalization\nfor Key Outcomes)" if lang == "en" else "LINKO\n(Latent Information Normalization\nfor Key Outcomes)\n\u4e3b\u8981\u30a2\u30a6\u30c8\u30ab\u30e0\u306e\u305f\u3081\u306e\u6f5c\u5728\u60c5\u5831\u6b63\u898f\u5316"
    _add_shape_box(slide, Inches(4.0), Inches(0.9), Inches(5.3), Inches(1.2),
                   linko_title, RGBColor(0x2C, 0x3E, 0x50), font_size=14, bold=True)
    # Make text white
    shape = slide.shapes[-1]
    for par in shape.text_frame.paragraphs:
        par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Core concept
    core_text = "Core Question: Does the endpoint carry the same\ninformational weight in every study?" if lang == "en" else "\u6838\u5fc3\u7684\u554f\u3044: \u30a8\u30f3\u30c9\u30dd\u30a4\u30f3\u30c8\u306f\u3059\u3079\u3066\u306e\n\u7814\u7a76\u3067\u540c\u3058\u60c5\u5831\u7684\u91cd\u307f\u3092\u6301\u3064\u304b\uff1f"
    _add_shape_box(slide, Inches(3.2), Inches(2.4), Inches(6.9), Inches(0.9),
                   core_text, RGBColor(0xEB, 0xF5, 0xFB), font_size=13)

    # Arrow down from core to two approaches
    _add_arrow(slide, Inches(5.0), Inches(3.3), Inches(3.5), Inches(3.8))
    _add_arrow(slide, Inches(8.3), Inches(3.3), Inches(9.8), Inches(3.8))

    # Left branch: Approach 1
    a1_title = "Approach 1: Variance-based ICR" if lang == "en" else "\u30a2\u30d7\u30ed\u30fc\u30c11: \u5206\u6563\u30d9\u30fc\u30b9ICR"
    _add_shape_box(slide, Inches(1.0), Inches(3.9), Inches(4.8), Inches(0.7),
                   a1_title, RGBColor(0x34, 0x98, 0xDB), font_size=13, bold=True)
    shape = slide.shapes[-1]
    for par in shape.text_frame.paragraphs:
        par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    a1_detail = "ICR_std = d / D\nComputable from published Table 1" if lang == "en" else "ICR_std = d / D\nTable 1\u304b\u3089\u8a08\u7b97\u53ef\u80fd"
    _add_shape_box(slide, Inches(1.0), Inches(4.8), Inches(4.8), Inches(0.7),
                   a1_detail, RGBColor(0xD6, 0xEA, 0xF8), font_size=11)

    # Right branch: Approach 2
    a2_title = "Approach 2: PCA-based ICR" if lang == "en" else "\u30a2\u30d7\u30ed\u30fc\u30c12: PCA\u30d9\u30fc\u30b9ICR"
    _add_shape_box(slide, Inches(7.5), Inches(3.9), Inches(4.8), Inches(0.7),
                   a2_title, RGBColor(0xE7, 0x4C, 0x3C), font_size=13, bold=True)
    shape = slide.shapes[-1]
    for par in shape.text_frame.paragraphs:
        par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    a2_detail = "Loading-based + Regression-based\nRequires individual patient data (IPD)" if lang == "en" else "Loading\u6cd5 + \u56de\u5e30\u6cd5\n\u500b\u7968\u30c7\u30fc\u30bf(IPD)\u304c\u5fc5\u8981"
    _add_shape_box(slide, Inches(7.5), Inches(4.8), Inches(4.8), Inches(0.7),
                   a2_detail, RGBColor(0xFA, 0xDB, 0xD8), font_size=11)

    # Arrows down to ICRD
    _add_arrow(slide, Inches(3.4), Inches(5.5), Inches(5.5), Inches(5.9))
    _add_arrow(slide, Inches(9.9), Inches(5.5), Inches(7.8), Inches(5.9))

    # Bottom: ICRD + output
    icrd_text = "ICR Discrepancy (ICRD) = max(ICR) - min(ICR)\n+ Prism Forest Plot Visualization" if lang == "en" else "ICR Discrepancy (ICRD) = max(ICR) - min(ICR)\n+ Prism Forest Plot\u53ef\u8996\u5316"
    _add_shape_box(slide, Inches(3.5), Inches(5.9), Inches(6.3), Inches(0.8),
                   icrd_text, RGBColor(0x27, 0xAE, 0x60), font_size=12, bold=True)
    shape = slide.shapes[-1]
    for par in shape.text_frame.paragraphs:
        par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Recommendation
    rec_text = "Report ICRD alongside I\u00b2 and \u03c4\u00b2 in meta-analyses" if lang == "en" else "\u30e1\u30bf\u89e3\u6790\u3067I\u00b2\u3001\u03c4\u00b2\u3068\u3068\u3082\u306bICRD\u3092\u5831\u544a"
    _add_shape_box(slide, Inches(3.5), Inches(6.85), Inches(6.3), Inches(0.5),
                   rec_text, RGBColor(0xF9, 0xE7, 0x9F), font_size=11, bold=True)

    return slide


def add_study_workflow_slide(prs, lang="en"):
    """Create editable Study Design / Analysis Workflow diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_text = "LINKO Study Design and Validation Workflow" if lang == "en" else "LINKO\u7814\u7a76\u30c7\u30b6\u30a4\u30f3\u3068\u691c\u8a3c\u30ef\u30fc\u30af\u30d5\u30ed\u30fc"
    tx = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Step boxes - horizontal flow
    steps = [
        ("1. Simulation\n(Monte Carlo)", "Scenario A: Uniform ICR\nScenario B: Heterogeneous ICR\nScenario C: Sequential", RGBColor(0x34, 0x98, 0xDB)),
        ("2. Real-World\n(Table 1 ICR)", "Statin therapy (5 RCTs)\nGlucose control (4 RCTs)", RGBColor(0xE6, 0x7E, 0x22)),
        ("3. PCA Validation\n(IST IPD)", "19,435 patients\n8 country sub-studies\n25 variables", RGBColor(0x27, 0xAE, 0x60)),
        ("4. Visualization\n(Prism Forest Plot)", "Color = ICR value\nSize = Secondary ICR\nSide panel = ICR bars", RGBColor(0x8E, 0x44, 0xAD)),
        ("5. Convergence\nAnalysis", "500 iterations\n3 strategies\nICR-guided selection", RGBColor(0xE7, 0x4C, 0x3C)),
    ] if lang == "en" else [
        ("1. \u30b7\u30df\u30e5\u30ec\u30fc\u30b7\u30e7\u30f3\n(\u30e2\u30f3\u30c6\u30ab\u30eb\u30ed)", "\u30b7\u30ca\u30eaA: \u5747\u4e00ICR\n\u30b7\u30ca\u30eaB: \u7570\u8cea ICR\n\u30b7\u30ca\u30eaC: \u9010\u6b21", RGBColor(0x34, 0x98, 0xDB)),
        ("2. \u5b9f\u30c7\u30fc\u30bf\n(Table 1 ICR)", "\u30b9\u30bf\u30c1\u30f3(5 RCT)\n\u8840\u7cd6\u30b3\u30f3\u30c8\u30ed\u30fc\u30eb(4 RCT)", RGBColor(0xE6, 0x7E, 0x22)),
        ("3. PCA\u691c\u8a3c\n(IST\u500b\u7968)", "19,435\u540d\n8\u304b\u56fd\u30b5\u30d6\u30b9\u30bf\u30c7\u30a3\n25\u5909\u6570", RGBColor(0x27, 0xAE, 0x60)),
        ("4. \u53ef\u8996\u5316\n(Prism Forest Plot)", "\u8272 = ICR\u5024\n\u30b5\u30a4\u30ba = \u526f\u6b21ICR\n\u30b5\u30a4\u30c9\u30d1\u30cd\u30eb = ICR\u68d2", RGBColor(0x8E, 0x44, 0xAD)),
        ("5. \u53ce\u675f\u5206\u6790", "500\u53cd\u5fa9\n3\u6226\u7565\nICR\u8a98\u5c0e\u9078\u629e", RGBColor(0xE7, 0x4C, 0x3C)),
    ]

    box_w = Inches(2.2)
    box_h_top = Inches(1.0)
    box_h_bot = Inches(1.5)
    gap = Inches(0.25)
    start_x = Inches(0.6)
    top_y = Inches(1.3)
    bot_y = Inches(2.6)

    for i, (header, detail, color) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        # Header box
        _add_shape_box(slide, x, top_y, box_w, box_h_top, header, color, font_size=13, bold=True)
        shape = slide.shapes[-1]
        for par in shape.text_frame.paragraphs:
            par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Detail box
        detail_color = RGBColor(min(255, color[0] + 80), min(255, color[1] + 80), min(255, color[2] + 80))
        _add_shape_box(slide, x, bot_y, box_w, box_h_bot, detail, detail_color, font_size=10)

        # Arrow between steps
        if i < len(steps) - 1:
            _add_arrow(slide, x + box_w, top_y + box_h_top / 2,
                       x + box_w + gap, top_y + box_h_top / 2)

    # Bottom section: Key results summary
    results_title = "Key Findings" if lang == "en" else "\u4e3b\u8981\u306a\u77e5\u898b"
    tx2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.4), Inches(12.0), Inches(0.5))
    tf2 = tx2.text_frame
    tf2.paragraphs[0].text = results_title
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.name = "Calibri"

    findings = [
        ("Simulation", "Heterogeneous ICR \u2192 Higher I\u00b2\n(11.7% vs 11.0%)", RGBColor(0x34, 0x98, 0xDB)),
        ("Real-World", "Low ICRD = Stable MA\nHigh ICRD = Heterogeneous MA", RGBColor(0xE6, 0x7E, 0x22)),
        ("PCA (IST)", "ICR_pca \u00d7 Mortality\nr = 0.90 (p = 0.003)", RGBColor(0x27, 0xAE, 0x60)),
        ("LOO Sensitivity", "Robust: r = 0.84\u20130.95\nAll p < 0.02", RGBColor(0x8E, 0x44, 0xAD)),
        ("Convergence", "LINKO: lowest mean\nstudies to I\u00b2<25%", RGBColor(0xE7, 0x4C, 0x3C)),
    ] if lang == "en" else [
        ("\u30b7\u30df\u30e5\u30ec\u30fc\u30b7\u30e7\u30f3", "\u7570\u8ceaICR \u2192 \u9ad8I\u00b2\n(11.7% vs 11.0%)", RGBColor(0x34, 0x98, 0xDB)),
        ("\u5b9f\u30c7\u30fc\u30bf", "\u4f4eICRD = \u5b89\u5b9aMA\n\u9ad8ICRD = \u7570\u8ceaMA", RGBColor(0xE6, 0x7E, 0x22)),
        ("PCA (IST)", "ICR_pca \u00d7 \u6b7b\u4ea1\u7387\nr = 0.90 (p = 0.003)", RGBColor(0x27, 0xAE, 0x60)),
        ("LOO\u611f\u5ea6\u5206\u6790", "\u9811\u5065: r = 0.84\u20130.95\n\u5168\u3066p < 0.02", RGBColor(0x8E, 0x44, 0xAD)),
        ("\u53ce\u675f\u5206\u6790", "LINKO: I\u00b2<25%\u307e\u3067\u306e\n\u5e73\u5747\u7814\u7a76\u6570\u6700\u5c11", RGBColor(0xE7, 0x4C, 0x3C)),
    ]

    for i, (label, detail, color) in enumerate(findings):
        x = start_x + i * (box_w + gap)
        _add_shape_box(slide, x, Inches(5.0), box_w, Inches(0.5), label, color, font_size=11, bold=True)
        shape = slide.shapes[-1]
        for par in shape.text_frame.paragraphs:
            par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        detail_color = RGBColor(min(255, color[0] + 80), min(255, color[1] + 80), min(255, color[2] + 80))
        _add_shape_box(slide, x, Inches(5.6), box_w, Inches(1.0), detail, detail_color, font_size=10)

    return slide


def add_icr_concept_slide(prs, lang="en"):
    """Create editable ICR Concept Diagram: Study A vs Study B comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_text = "ICR Concept: Why Data Dimensionality Matters" if lang == "en" else "ICR\u6982\u5ff5: \u306a\u305c\u30c7\u30fc\u30bf\u6b21\u5143\u304c\u91cd\u8981\u304b"
    tx = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Study A: D=10
    sa_label = "Study A (D = 10)" if lang == "en" else "\u7814\u7a76A (D = 10)"
    _add_shape_box(slide, Inches(1.0), Inches(1.2), Inches(4.5), Inches(0.6),
                   sa_label, RGBColor(0x2C, 0x3E, 0x50), font_size=16, bold=True)
    shape = slide.shapes[-1]
    for par in shape.text_frame.paragraphs:
        par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Study A grid: 10 small boxes, 1 highlighted
    for j in range(10):
        x = Inches(1.0) + j * Inches(0.45)
        color = RGBColor(0xE7, 0x4C, 0x3C) if j == 0 else RGBColor(0xBD, 0xC3, 0xC7)
        _add_shape_box(slide, x, Inches(2.0), Inches(0.4), Inches(0.8),
                       "EP" if j == 0 else f"V{j}", color, font_size=9)
        if j == 0:
            shape = slide.shapes[-1]
            for par in shape.text_frame.paragraphs:
                par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                par.font.bold = True

    icr_a_text = "ICR = 1/10 = 0.10 (10%)" if lang == "en" else "ICR = 1/10 = 0.10 (10%)"
    _add_shape_box(slide, Inches(1.0), Inches(3.0), Inches(4.5), Inches(0.5),
                   icr_a_text, RGBColor(0xFA, 0xDB, 0xD8), font_size=14, bold=True)

    # Study B: D=50
    sb_label = "Study B (D = 50)" if lang == "en" else "\u7814\u7a76B (D = 50)"
    _add_shape_box(slide, Inches(1.0), Inches(4.0), Inches(10.5), Inches(0.6),
                   sb_label, RGBColor(0x2C, 0x3E, 0x50), font_size=16, bold=True)
    shape = slide.shapes[-1]
    for par in shape.text_frame.paragraphs:
        par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Study B grid: show ~25 boxes (representing 50)
    n_show = 25
    for j in range(n_show):
        x = Inches(1.0) + j * Inches(0.42)
        color = RGBColor(0xE7, 0x4C, 0x3C) if j == 0 else RGBColor(0xBD, 0xC3, 0xC7)
        label = "EP" if j == 0 else ("..." if j == n_show - 1 else "")
        _add_shape_box(slide, x, Inches(4.8), Inches(0.37), Inches(0.5),
                       label, color, font_size=7)
        if j == 0:
            shape = slide.shapes[-1]
            for par in shape.text_frame.paragraphs:
                par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                par.font.bold = True

    icr_b_text = "ICR = 1/50 = 0.02 (2%)" if lang == "en" else "ICR = 1/50 = 0.02 (2%)"
    _add_shape_box(slide, Inches(1.0), Inches(5.5), Inches(10.5), Inches(0.5),
                   icr_b_text, RGBColor(0xD6, 0xEA, 0xF8), font_size=14, bold=True)

    # Bottom question
    q_text = ("Same endpoint, same effect size estimate\n"
              "\u2192 But endpoint represents 10% vs 2% of the data space\n"
              "\u2192 ICRD = 0.08 \u2192 Are these truly comparable for pooling?") if lang == "en" else \
             ("\u540c\u3058\u30a8\u30f3\u30c9\u30dd\u30a4\u30f3\u30c8\u3001\u540c\u3058\u52b9\u679c\u91cf\u63a8\u5b9a\u5024\n"
              "\u2192 \u3057\u304b\u3057EP\u306f\u30c7\u30fc\u30bf\u7a7a\u9593\u306e10% vs 2%\n"
              "\u2192 ICRD = 0.08 \u2192 \u30d7\u30fc\u30ea\u30f3\u30b0\u306f\u59a5\u5f53\u304b\uff1f")
    _add_shape_box(slide, Inches(2.5), Inches(6.2), Inches(8.3), Inches(1.1),
                   q_text, RGBColor(0xF9, 0xE7, 0x9F), font_size=13, bold=True)

    return slide


def add_prism_concept_slide(prs, lang="en"):
    """Create editable Prism Forest Plot concept diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_text = "Prism Forest Plot: Concept" if lang == "en" else "Prism Forest Plot: \u6982\u5ff5"
    tx = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Left: Standard Forest Plot
    std_label = "Standard Forest Plot" if lang == "en" else "\u6a19\u6e96\u30d5\u30a9\u30ec\u30b9\u30c8\u30d7\u30ed\u30c3\u30c8"
    _add_shape_box(slide, Inches(0.5), Inches(1.0), Inches(3.5), Inches(0.6),
                   std_label, RGBColor(0x7F, 0x8C, 0x8D), font_size=14, bold=True)
    shape = slide.shapes[-1]
    for par in shape.text_frame.paragraphs:
        par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    std_items = [
        ("x-axis: Effect size", "\u2022"),
        ("Study weight: Inverse variance", "\u2022"),
        ("All bars same color", "\u2022"),
        ("No ICR information visible", "\u2717"),
    ] if lang == "en" else [
        ("x\u8ef8: \u52b9\u679c\u91cf", "\u2022"),
        ("\u91cd\u307f: \u9006\u5206\u6563", "\u2022"),
        ("\u5168\u30d0\u30fc\u540c\u8272", "\u2022"),
        ("ICR\u60c5\u5831\u306a\u3057", "\u2717"),
    ]
    for i, (text, bullet) in enumerate(std_items):
        y = Inches(1.8) + i * Inches(0.5)
        tx2 = slide.shapes.add_textbox(Inches(0.7), y, Inches(3.3), Inches(0.45))
        tf2 = tx2.text_frame
        tf2.paragraphs[0].text = f"{bullet} {text}"
        tf2.paragraphs[0].font.size = Pt(12)
        tf2.paragraphs[0].font.name = "Calibri"

    # Arrow in the middle
    _add_shape_box(slide, Inches(4.5), Inches(2.5), Inches(1.5), Inches(0.7),
                   "Prism\n\u2192", RGBColor(0x8E, 0x44, 0xAD), font_size=16, bold=True)
    shape = slide.shapes[-1]
    for par in shape.text_frame.paragraphs:
        par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Right: Prism Forest Plot
    prism_label = "LINKO Prism Forest Plot" if lang == "en" else "LINKO Prism Forest Plot"
    _add_shape_box(slide, Inches(6.5), Inches(1.0), Inches(6.0), Inches(0.6),
                   prism_label, RGBColor(0x8E, 0x44, 0xAD), font_size=14, bold=True)
    shape = slide.shapes[-1]
    for par in shape.text_frame.paragraphs:
        par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prism_items = [
        ("x-axis: Effect size (unchanged)", RGBColor(0x27, 0xAE, 0x60)),
        ("BAR COLOR = ICR value\n  (warm red = high, cool blue = low)", RGBColor(0xE7, 0x4C, 0x3C)),
        ("POINT SIZE = Secondary ICR\n  (e.g., ICR_pca regression)", RGBColor(0xE6, 0x7E, 0x22)),
        ("SIDE PANEL = ICR bar chart\n  aligned with each study", RGBColor(0x34, 0x98, 0xDB)),
    ] if lang == "en" else [
        ("x\u8ef8: \u52b9\u679c\u91cf (\u5909\u66f4\u306a\u3057)", RGBColor(0x27, 0xAE, 0x60)),
        ("\u30d0\u30fc\u306e\u8272 = ICR\u5024\n  (\u6696\u8272=\u9ad8, \u5bd2\u8272=\u4f4e)", RGBColor(0xE7, 0x4C, 0x3C)),
        ("\u70b9\u306e\u30b5\u30a4\u30ba = \u526f\u6b21ICR\n  (\u4f8b: ICR_pca\u56de\u5e30)", RGBColor(0xE6, 0x7E, 0x22)),
        ("\u30b5\u30a4\u30c9\u30d1\u30cd\u30eb = ICR\u68d2\u30b0\u30e9\u30d5\n  \u5404\u7814\u7a76\u3068\u6574\u5217", RGBColor(0x34, 0x98, 0xDB)),
    ]
    for i, (text, color) in enumerate(prism_items):
        y = Inches(1.8) + i * Inches(0.85)
        _add_shape_box(slide, Inches(6.5), y, Inches(6.0), Inches(0.75),
                       text, color, font_size=11, align=PP_ALIGN.LEFT)
        shape = slide.shapes[-1]
        for par in shape.text_frame.paragraphs:
            par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Bottom: analogy
    analogy = ('Like a prism decomposes white light into its spectrum,\n'
               'the Prism Forest Plot decomposes the forest plot to reveal\n'
               'hidden ICR dimensions of each study.') if lang == "en" else \
              ('\u30d7\u30ea\u30ba\u30e0\u304c\u767d\u8272\u5149\u3092\u30b9\u30da\u30af\u30c8\u30eb\u306b\u5206\u89e3\u3059\u308b\u3088\u3046\u306b\u3001\n'
               'Prism Forest Plot\u306f\u30d5\u30a9\u30ec\u30b9\u30c8\u30d7\u30ed\u30c3\u30c8\u3092\u5206\u89e3\u3057\u3066\n'
               '\u5404\u7814\u7a76\u306e\u96a0\u308c\u305fICR\u6b21\u5143\u3092\u660e\u3089\u304b\u306b\u3059\u308b\u3002')
    _add_shape_box(slide, Inches(2.0), Inches(5.8), Inches(9.3), Inches(1.0),
                   analogy, RGBColor(0xF5, 0xF5, 0xF5), font_size=13, bold=False)

    return slide


def add_table_slide(prs, title_text, headers, rows, caption=""):
    """Add a slide with an editable table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    tx = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Table
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = min(Inches(12.0), Inches(n_cols * 2.0))
    left = int((SLIDE_W - tbl_width) / 2)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, Inches(1.0), tbl_width, Inches(0.4 * n_rows))
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for par in cell.text_frame.paragraphs:
            par.font.size = Pt(11)
            par.font.bold = True
            par.font.name = "Calibri"
            par.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        for par in cell.text_frame.paragraphs:
            par.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for par in cell.text_frame.paragraphs:
                par.font.size = Pt(10)
                par.font.name = "Calibri"
                par.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF3, 0xF4)

    # Caption
    if caption:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2 + 0.4 * n_rows), Inches(12.3), Inches(0.4))
        tf2 = tx2.text_frame
        tf2.word_wrap = True
        tf2.paragraphs[0].text = caption
        tf2.paragraphs[0].font.size = Pt(10)
        tf2.paragraphs[0].font.italic = True
        tf2.paragraphs[0].font.name = "Calibri"
        tf2.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    return slide


# ============================================================
# BUILD PRESENTATIONS
# ============================================================
def build_pptx(lang="en"):
    prs = new_pres()

    # --- Slide 1: Title ---
    if lang == "en":
        add_title_slide(prs,
            "LINKO (Latent Information Normalization for Key Outcomes)",
            "A Framework for Evaluating the Validity of\nMeta-Analytic Pooling Across Heterogeneous RCT Data Structures")
    else:
        add_title_slide(prs,
            "LINKO (Latent Information Normalization for Key Outcomes)",
            "RCT\u30c7\u30fc\u30bf\u69cb\u9020\u306e\u7570\u8cea\u6027\u3092\u8003\u616e\u3057\u305f\n\u30e1\u30bf\u89e3\u6790\u30d7\u30fc\u30ea\u30f3\u30b0\u306e\u59a5\u5f53\u6027\u8a55\u4fa1\u30d5\u30ec\u30fc\u30e0\u30ef\u30fc\u30af")

    # --- Slide 2: ICR Concept (editable) ---
    add_icr_concept_slide(prs, lang)

    # --- Slide 3: LINKO Framework Overview (editable) ---
    add_linko_framework_slide(prs, lang)

    # --- Slide 4: Study Workflow (editable) ---
    add_study_workflow_slide(prs, lang)

    # --- Slide 5: Prism Forest Plot concept (editable) ---
    add_prism_concept_slide(prs, lang)

    # --- Slide 6: Fig 0 - ICR Dimension Relationship (code output) ---
    fig_title = "Figure 1. Theoretical ICR as a Function of Data Dimensionality" if lang == "en" else "\u56f31. \u30c7\u30fc\u30bf\u6b21\u5143\u306e\u95a2\u6570\u3068\u3057\u3066\u306e\u7406\u8ad6\u7684ICR"
    add_image_slide(prs, "fig0_icr_dimension_relationship.png", fig_title)

    # --- Slide 7: Fig 1 - Scenario Comparison (code output) ---
    fig_title = "Figure 2. Simulation: Scenario A vs B Comparison" if lang == "en" else "\u56f32. \u30b7\u30df\u30e5\u30ec\u30fc\u30b7\u30e7\u30f3: \u30b7\u30ca\u30eaA\u5bfeB\u306e\u6bd4\u8f03"
    add_image_slide(prs, "fig1_scenario_comparison.png", fig_title)

    # --- Slide 8: Fig 2 - Sequential Analysis (code output) ---
    fig_title = "Figure 3. Sequential Meta-Analysis (Scenario C)" if lang == "en" else "\u56f33. \u9010\u6b21\u30e1\u30bf\u89e3\u6790 (\u30b7\u30ca\u30eaC)"
    add_image_slide(prs, "fig2_sequential_analysis.png", fig_title)

    # --- Slide 9: Table 1 - Simulation Results ---
    t1_title = "Table 1. Simulation Results" if lang == "en" else "\u88681. \u30b7\u30df\u30e5\u30ec\u30fc\u30b7\u30e7\u30f3\u7d50\u679c"
    if lang == "en":
        add_table_slide(prs, t1_title,
            ["Metric", "Scenario A (Uniform ICR)", "Scenario B (Heterogeneous ICR)"],
            [["Mean I\u00b2 (SD)", "11.0% (17.2)", "11.7% (16.0)"],
             ["Mean ICRD (SD)", "0.0000", "0.1753 (0.031)"],
             ["Mean Pooled Effect", "0.500", "0.496"],
             ["Effect Bias", "+0.0004", "-0.0042"]])
    else:
        add_table_slide(prs, t1_title,
            ["\u6307\u6a19", "\u30b7\u30ca\u30eaA (\u5747\u4e00ICR)", "\u30b7\u30ca\u30eaB (\u7570\u8ceaICR)"],
            [["\u5e73\u5747I\u00b2 (SD)", "11.0% (17.2)", "11.7% (16.0)"],
             ["\u5e73\u5747ICRD (SD)", "0.0000", "0.1753 (0.031)"],
             ["\u5e73\u5747\u30d7\u30fc\u30eb\u52b9\u679c", "0.500", "0.496"],
             ["\u52b9\u679c\u30d0\u30a4\u30a2\u30b9", "+0.0004", "-0.0042"]])

    # --- Slide 10: Real-world Statin (code output) ---
    fig_title = "Figure 4. ICR Analysis: Statin Therapy" if lang == "en" else "\u56f34. ICR\u5206\u6790: \u30b9\u30bf\u30c1\u30f3\u7642\u6cd5"
    add_image_slide(prs, "fig_realworld_statin.png", fig_title)

    # --- Slide 11: Real-world Glucose (code output) ---
    fig_title = "Figure 5. ICR Analysis: Intensive Glucose Control" if lang == "en" else "\u56f35. ICR\u5206\u6790: \u5f37\u5316\u8840\u7cd6\u30b3\u30f3\u30c8\u30ed\u30fc\u30eb"
    add_image_slide(prs, "fig_realworld_glucose_control.png", fig_title)

    # --- Slide 12: Table 2 & 3 - Real-world ICR ---
    t2_title = "Table 2. Statin Therapy RCTs" if lang == "en" else "\u88682. \u30b9\u30bf\u30c1\u30f3\u7642\u6cd5RCT"
    add_table_slide(prs, t2_title,
        ["Study", "D", "d", "ICR_std", "Effect Size"],
        [["4S (1994)", "10", "1", "0.100", "-0.370"],
         ["WOSCOPS (1995)", "10", "1", "0.100", "-0.250"],
         ["CARE (1996)", "11", "1", "0.091", "-0.100"],
         ["LIPID (1998)", "10", "1", "0.100", "-0.250"],
         ["AFCAPS (1998)", "10", "1", "0.100", "-0.110"]],
        "ICRD = 0.009, I\u00b2 = 0.0%, Pooled = -0.251")

    t3_title = "Table 3. Glucose Control RCTs" if lang == "en" else "\u88683. \u8840\u7cd6\u30b3\u30f3\u30c8\u30ed\u30fc\u30ebRCT"
    add_table_slide(prs, t3_title,
        ["Study", "D", "d", "ICR_std", "Effect Size"],
        [["UKPDS 33 (1998)", "8", "1", "0.125", "-0.060"],
         ["ACCORD (2008)", "13", "1", "0.077", "+0.220"],
         ["ADVANCE (2008)", "12", "1", "0.083", "-0.070"],
         ["VADT (2009)", "12", "1", "0.083", "-0.020"]],
        "ICRD = 0.048, I\u00b2 = 17.0%, Pooled = -0.003")

    # --- Slide 14: PCA IST Analysis (code output) ---
    fig_title = "Figure 6. PCA-Based ICR Analysis (IST)" if lang == "en" else "\u56f36. PCA\u30d9\u30fc\u30b9ICR\u5206\u6790 (IST)"
    add_image_slide(prs, "fig_pca_ist_analysis.png", fig_title)

    # --- Slide 15: Table 4 - IST PCA ---
    t4_title = "Table 4. IST PCA-Based ICR by Country" if lang == "en" else "\u88684. IST PCA\u30d9\u30fc\u30b9ICR (\u56fd\u5225)"
    add_table_slide(prs, t4_title,
        ["Country", "n", "Mortality", "ICR_std", "ICR_pca (loading)", "ICR_pca (reg)", "PCs"],
        [["UK", "5,787", "28.6%", "0.040", "0.138", "0.00162", "4"],
         ["Italy", "3,112", "20.0%", "0.040", "0.046", "0.00153", "2"],
         ["Switzerland", "1,631", "23.1%", "0.040", "0.121", "0.00135", "4"],
         ["Poland", "759", "29.6%", "0.040", "0.139", "0.00230", "4"],
         ["Netherlands", "728", "18.3%", "0.040", "0.180", "0.00136", "5"],
         ["Sweden", "636", "12.7%", "0.040", "0.096", "0.00073", "3"],
         ["Australia", "568", "15.0%", "0.040", "0.109", "0.00096", "4"],
         ["Argentina", "545", "21.8%", "0.040", "0.079", "0.00157", "3"]],
        "r(ICR_pca_reg, mortality) = 0.90, p = 0.003")

    # --- Slide 16: LOO Sensitivity (code output) ---
    fig_title = "Figure 7. Leave-One-Out Sensitivity Analysis" if lang == "en" else "\u56f37. Leave-One-Out\u611f\u5ea6\u5206\u6790"
    add_image_slide(prs, "fig_loo_sensitivity.png", fig_title)

    # --- Slide 17: Table 5 - LOO ---
    t5_title = "Table 5. Leave-One-Out Sensitivity" if lang == "en" else "\u88685. Leave-One-Out\u611f\u5ea6\u5206\u6790"
    add_table_slide(prs, t5_title,
        ["Excluded", "r (loading)", "p", "r (regression)", "p"],
        [["UK", "0.169", "0.717", "0.954", "0.001"],
         ["Italy", "0.286", "0.534", "0.908", "0.005"],
         ["Switzerland", "0.258", "0.577", "0.914", "0.004"],
         ["Poland", "0.154", "0.742", "0.860", "0.013"],
         ["Netherlands", "0.526", "0.225", "0.903", "0.005"],
         ["Sweden", "0.205", "0.659", "0.843", "0.017"],
         ["Australia", "0.272", "0.555", "0.875", "0.010"],
         ["Argentina", "0.299", "0.515", "0.898", "0.006"],
         ["Full (n=8)", "0.265", "0.526", "0.896", "0.003"]],
        "Regression-based ICR_pca robust: r = 0.84\u20130.95, all p < 0.02")

    # --- Slide 18: Prism Forest Plot - Statin (code output) ---
    fig_title = "Figure 8. LINKO Prism Forest Plot: Statin Therapy" if lang == "en" else "\u56f38. Prism Forest Plot: \u30b9\u30bf\u30c1\u30f3\u7642\u6cd5"
    cap = "Low ICRD = 0.009: uniform color across all studies" if lang == "en" else "\u4f4eICRD = 0.009: \u5168\u7814\u7a76\u3067\u5747\u4e00\u306a\u8272"
    add_image_slide(prs, "fig_linko_prism_statin.png", fig_title, cap)

    # --- Slide 19: Prism Forest Plot - Glucose (code output) ---
    fig_title = "Figure 9. LINKO Prism Forest Plot: Glucose Control" if lang == "en" else "\u56f39. Prism Forest Plot: \u8840\u7cd6\u30b3\u30f3\u30c8\u30ed\u30fc\u30eb"
    cap = "High ICRD = 0.048: visible color gradient (UKPDS warm vs ACCORD cool)" if lang == "en" else "\u9ad8ICRD = 0.048: \u8272\u30b0\u30e9\u30c7\u30a3\u30a8\u30f3\u30c8\u304c\u53ef\u8996\u7684"
    add_image_slide(prs, "fig_linko_prism_glucose.png", fig_title, cap)

    # --- Slide 20: Prism Forest Plot - IST (code output) ---
    fig_title = "Figure 10. LINKO Prism Forest Plot: IST Country Sub-Studies" if lang == "en" else "\u56f310. Prism Forest Plot: IST\u56fd\u5225\u30b5\u30d6\u30b9\u30bf\u30c7\u30a3"
    cap = "Color = ICR_pca (loading), Size = ICR_pca (regression)" if lang == "en" else "\u8272 = ICR_pca (loading\u6cd5), \u30b5\u30a4\u30ba = ICR_pca (\u56de\u5e30\u6cd5)"
    add_image_slide(prs, "fig_linko_prism_ist.png", fig_title, cap)

    # --- Slide 21: Early Convergence (code output) ---
    fig_title = "Figure 11. LINKO Early Convergence Analysis" if lang == "en" else "\u56f311. LINKO\u65e9\u671f\u53ce\u675f\u5206\u6790"
    add_image_slide(prs, "fig_linko_early_convergence.png", fig_title)

    # --- Slide 22: Table 6 - Early Convergence ---
    t6_title = "Table 6. Early Convergence Results (500 iterations)" if lang == "en" else "\u88686. \u65e9\u671f\u53ce\u675f\u7d50\u679c (500\u53cd\u5fa9)"
    if lang == "en":
        add_table_slide(prs, t6_title,
            ["Strategy", "Mean to conclusion", "Median", "% by 5", "% by 10", "Mean to I\u00b2<25%"],
            [["Random", "3.94", "3.0", "81.4%", "97.2%", "3.19"],
             ["ICR-matched", "3.91", "3.0", "78.0%", "97.6%", "3.20"],
             ["LINKO optimized", "4.00", "3.0", "77.6%", "97.6%", "3.14"]],
            "LINKO achieved lowest mean studies for I\u00b2 < 25% (3.14 vs 3.19\u20133.20)")
    else:
        add_table_slide(prs, t6_title,
            ["\u6226\u7565", "\u7d50\u8ad6\u5e73\u5747", "\u4e2d\u592e\u5024", "5\u4ee5\u5185%", "10\u4ee5\u5185%", "I\u00b2<25%\u5e73\u5747"],
            [["\u30e9\u30f3\u30c0\u30e0", "3.94", "3.0", "81.4%", "97.2%", "3.19"],
             ["ICR\u30de\u30c3\u30c1", "3.91", "3.0", "78.0%", "97.6%", "3.20"],
             ["LINKO\u6700\u9069\u5316", "4.00", "3.0", "77.6%", "97.6%", "3.14"]],
            "LINKO: I\u00b2<25%\u307e\u3067\u306e\u5e73\u5747\u7814\u7a76\u6570\u6700\u5c11 (3.14 vs 3.19\u20133.20)")

    # --- Slide 23: Summary comparison ---
    t7_title = "Table 7. Summary: ICR Approaches Comparison" if lang == "en" else "\u88687. \u307e\u3068\u3081: ICR\u30a2\u30d7\u30ed\u30fc\u30c1\u306e\u6bd4\u8f03"
    if lang == "en":
        add_table_slide(prs, t7_title,
            ["Analysis", "ICR Measure", "Variation (CV)", "Association"],
            [["Statin therapy", "ICR_std (Approach 1)", "0.041", "I\u00b2 = 0.0% (stable)"],
             ["Glucose control", "ICR_std (Approach 1)", "0.240", "I\u00b2 = 17.0% (heterogeneous)"],
             ["IST (loading)", "ICR_pca (Approach 2)", "0.360", "r = 0.27 with mortality"],
             ["IST (regression)", "ICR_pca_reg (Approach 2)", "0.329", "r = 0.90 with mortality"]])
    else:
        add_table_slide(prs, t7_title,
            ["\u5206\u6790", "ICR\u6307\u6a19", "\u5909\u52d5 (CV)", "\u95a2\u9023"],
            [["\u30b9\u30bf\u30c1\u30f3", "ICR_std (\u30a2\u30d7\u30ed\u30fc\u30c11)", "0.041", "I\u00b2 = 0.0% (\u5b89\u5b9a)"],
             ["\u8840\u7cd6\u30b3\u30f3\u30c8\u30ed\u30fc\u30eb", "ICR_std (\u30a2\u30d7\u30ed\u30fc\u30c11)", "0.240", "I\u00b2 = 17.0% (\u7570\u8cea)"],
             ["IST (loading)", "ICR_pca (\u30a2\u30d7\u30ed\u30fc\u30c12)", "0.360", "\u6b7b\u4ea1\u7387\u3068r = 0.27"],
             ["IST (\u56de\u5e30)", "ICR_pca_reg (\u30a2\u30d7\u30ed\u30fc\u30c12)", "0.329", "\u6b7b\u4ea1\u7387\u3068r = 0.90"]])

    # Save
    suffix = "english" if lang == "en" else "japanese"
    out = os.path.join(BASE, f"ICR_figures_{suffix}.pptx")
    prs.save(out)
    print(f"Saved {suffix.upper()} PPTX: {out}")
    return out


if __name__ == "__main__":
    print("Generating English PPTX...")
    build_pptx("en")
    print("Generating Japanese PPTX...")
    build_pptx("ja")
    print("Done.")
